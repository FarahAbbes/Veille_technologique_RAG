from typing import List, Any, Optional, Dict


class ImageExtractor:
    def extract(self, pdf_path: str) -> List[Any]:
        try:
            import fitz
        except Exception:
            return []
        out = []
        doc = fitz.open(pdf_path)
        for i in range(len(doc)):
            for img in doc.get_page_images(i):
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                out.append(pix)
        return out


class OCRProcessor:
    def ocr_image(self, image) -> str:
        try:
            import pytesseract
            from PIL import Image
        except Exception:
            return ""
        try:
            if not isinstance(image, Image.Image):
                return ""
            return pytesseract.image_to_string(image)
        except Exception:
            return ""


class TableExtractor:
    def extract(self, pdf_path: str) -> List[Dict[str, Any]]:
        try:
            import camelot
            tables = camelot.read_pdf(pdf_path, pages="all")
            return [{"page": t.page, "shape": t.shape} for t in tables]
        except Exception:
            pass
        try:
            import tabula
            dfs = tabula.read_pdf(pdf_path, pages="all")
            return [{"rows": len(df), "cols": len(df.columns)} for df in dfs]
        except Exception:
            return []


class MultiFormatExtractor:
    def extract_docx(self, path: str) -> str:
        try:
            import docx
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""

    def extract_pptx(self, path: str) -> List[str]:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return texts
        except Exception:
            return []

    def extract_xlsx(self, path: str) -> List[Dict[str, Any]]:
        try:
            import pandas as pd
            xls = pd.ExcelFile(path)
            out = []
            for sheet in xls.sheet_names:
                df = xls.parse(sheet)
                out.append({"sheet": sheet, "rows": len(df), "cols": len(df.columns)})
            return out
        except Exception:
            return []
